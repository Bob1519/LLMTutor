from pptx import Presentation


#extracts the text from a powerpoint and places it into a text file
file_name = "Hypothesis testing part 5_posted"
pres = Presentation(f"{file_name}.pptx")

f = open(f"{file_name}.txt", "x")

for i, slide in enumerate(pres.slides):
    print(f"Slide {i+1}")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
            f.write(shape.text + '\n')